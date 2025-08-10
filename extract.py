import os
from pptx import Presentation
from PIL import Image
import pytesseract
from concurrent.futures import ThreadPoolExecutor, as_completed
from utils import file_hash, cache_load, cache_save

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'


def extract_text_from_pptx(pptx_path):
    cache_key = f"pptx_{file_hash(pptx_path)}"
    cached = cache_load(cache_key)
    if cached:
        print(f"[CACHE] Using cached text for {pptx_path}")
        return cached

    prs = Presentation(pptx_path)
    slides_text = {}
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)
        slides_text[i] = " ".join(texts)

    cache_save(cache_key, slides_text)
    return slides_text


def ocr_image(img_path):
    try:
        text = pytesseract.image_to_string(Image.open(img_path))
        return img_path, text.strip()
    except Exception as e:
        return img_path, f"[ERROR] OCR failed: {str(e)}"


def extract_text_from_images(folder_path):
    image_files = sorted([f for f in os.listdir(folder_path) 
                         if f.lower().endswith(('.png', '.jpg', '.jpeg'))])
    cache_key = f"imgfolder_{file_hash(','.join(image_files))}"
    
    cached = cache_load(cache_key)
    if cached:
        print(f"[CACHE] Using cached OCR for {folder_path}")
        return cached

    slides_text = {}
    with ThreadPoolExecutor() as executor:
        futures = {executor.submit(ocr_image, os.path.join(folder_path, img)): idx + 1 
                  for idx, img in enumerate(image_files)}
        
        for future in as_completed(futures):
            slide_num = futures[future]
            _, text = future.result()
            slides_text[slide_num] = text

    cache_save(cache_key, slides_text)
    return slides_text
