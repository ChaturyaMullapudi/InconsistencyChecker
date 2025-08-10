import os
import json
import argparse
import requests
from pptx import Presentation
from collections import defaultdict
from PIL import Image
import pytesseract
import dotenv

dotenv.load_dotenv()

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# ====== CONFIG ======
GEMINI_API_KEY = os.getenv("GOOGLE_API_KEY")
GEMINI_MODEL = "gemini-2.5-flash"
GEMINI_ENDPOINT = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent"

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slides_text = {}
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)
        slides_text[i] = " ".join(texts)
    return slides_text

def extract_text_from_images(folder_path):
    slides_text = {}
    image_files = sorted([f for f in os.listdir(folder_path) if f.lower().endswith(('.png', '.jpg', '.jpeg'))])
    for i, img_file in enumerate(image_files, start=1):
        img_path = os.path.join(folder_path, img_file)
        text = pytesseract.image_to_string(Image.open(img_path))
        slides_text[i] = text.strip()
    return slides_text


def run_gemini_analysis(slides_text):
    prompt = f"""
            You are an expert at reviewing PowerPoint decks for factual and logical inconsistencies.
            Include  Textual contradictions and contextual Logic gaps also. Look Into EVERY SLIDE. DONT MISS ANY

            Here is the extracted slide text:
            {json.dumps(slides_text, indent=2)}


            Instructions:
            1. Validate the local findings and include them in the final output, even if you find them correct as-is.
            2. Identify any additional contradictions or inconsistencies in the data, dates, or claims.
            3. Clearly specify the slides involved in each issue.
            4. Also include if any content deviates from its primary claim or context.
            5. Output ALL results in this format:

            [Slide X & Slide Y] ISSUE: <description>
            """

    headers = {"Content-Type": "application/json"}
    params = {"key": GEMINI_API_KEY}
    body = {
        "contents": [{
            "parts": [{"text": prompt}]
        }]
    }
    resp = requests.post(GEMINI_ENDPOINT, headers=headers, params=params, json=body)
    resp.raise_for_status()
    data = resp.json()
    if "candidates" in data and data["candidates"]:
        return data["candidates"][0]["content"]["parts"][0]["text"]
    else:
        return "No AI findings returned."


def main():
    parser = argparse.ArgumentParser(description="Noogat Inconsistency Checker")
    parser.add_argument("--pptx", help="Path to PowerPoint file")
    parser.add_argument("--images", help="Path to image file")
    args = parser.parse_args()

    if args.pptx:
        slides_text = extract_text_from_pptx(args.pptx)
    elif args.images:
        slides_text = extract_text_from_images(args.images)
    else:
        parser.error("Either --pptx or --images argument is required")
    

    # Running our LLM (gemini)
    gemini_findings = run_gemini_analysis(slides_text)

    print("\n===== GEMINI AI FINDINGS =====")
    print(gemini_findings)


if __name__ == "__main__":
    main()
